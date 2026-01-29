import datetime
import itertools
import argparse
import os
import numbers
import math

import requests
import pandas
import pptx
import pptx.chart.data
import yaml

# Setting variables
V1_TOKEN = os.environ.get('TMV1_TOKEN', '')
# Specify the correct domain name for your region in V1_URL
#   default: https://api.xdr.trendmicro.com (US region)
V1_URL = os.environ.get('TMV1_URL', 'https://api.xdr.trendmicro.com')
# This value is used for User-Agent header in API requests. So you can
# customize this value to describe your company name, integration tool name,
# and so on as you like.
#   default: "Trend Vision One API Cookbook ({script_name})"
V1_UA = os.environ.get('TMV1_UA', 'Trend Vision One API Cookbook '
                       f'({os.path.basename(__file__)})')
V1_XLSX_FILENAME = os.environ.get('TMV1_XLSX_FILENAME',
                                  'security_posture.xlsx')
V1_PPTX_FILENAME = os.environ.get('TMV1_PPTX_FILENAME',
                                  'security_posture.pptx')
V1_YAML_FILENAME = os.environ.get(
    'TMV1_YAML_FILENAME',
    os.path.join(
        os.path.dirname(os.path.realpath(__file__)), 'security_posture.yaml'
    )
)


def is_container(v):
    try:
        if isinstance(v, (str, bytes)):
            return False
        iter(v)
    except TypeError:
        return False
    return True


def split_camel_case(s):
    if not s:
        return []
    indexes = [0]
    length = len(s)
    indexes.extend(i for i in range(1, length-1) if (
        s[i].isupper() and (s[i-1].islower() or s[i+1].islower())
    ))
    indexes.append(length)
    return [s[indexes[i]:indexes[i+1]] for i in range(len(indexes) - 1)]


def unique_list(v):
    return list(dict.fromkeys(v))


class TmV1Client:
    base_url_default = V1_URL

    def __init__(self, token, base_url=None):
        if not token:
            raise ValueError('Authentication token missing')
        self.token = token
        self.base_url = base_url or TmV1Client.base_url_default

    def make_headers(self, **kwargs):
        headers = {}
        use_token = kwargs.pop('use_token', True)
        if use_token:
            headers['Authorization'] = 'Bearer ' + self.token
        headers['User-Agent'] = V1_UA
        return headers

    def get(self, url_or_path, use_token=True, **kwargs):
        kwargs.setdefault('headers', {}).update(
            self.make_headers(use_token=use_token)
        )
        url = (self.base_url + url_or_path if url_or_path.startswith('/')
               else url_or_path)
        r = requests.get(url, **kwargs)
        if 200 == r.status_code:
            if 'application/json' in r.headers.get('Content-Type', ''):
                return r.json()
            return r.content
        raise RuntimeError(f'Request unsuccessful (GET {url_or_path}):'
                           f' {r.status_code} {r.text}')

    def get_security_posture(self):
        return self.get('/v3.0/asrm/securityPosture')


unit_suffixes = ['count', 'density', 'days', 'rate']


def get_unit(s):
    suffix = split_camel_case(s)[-1].lower()
    return suffix if suffix in unit_suffixes else ''


def iter_nested_dict(value, keys=[]):
    """
    This yields nested dicts that don't contain any dict value.
    After yields dicts, this removes them from their parent dict.
    So, this also yields these parents as a dict without any dict values.
    If the units (suffix) of all keys are the same, this yields the parent
    dict. Or, this yield a dict containing only each key-value pair.

    When handling a list including dicts, this replaces the list-index of
    each dict to the dict-key with the value of its 1st key.
    After that, if all dicts in a list have only one value for the same key,
    this combine them to one dict.
    When a list including dicts is empty, this skips to record it.
    """
    if isinstance(value, dict):
        units = set()
        iters = {}
        key_copy = list(value.keys())
        for k in key_copy:
            if keys and isinstance(keys[-1], int):
                # When parent is list, replace index to value of the 1st key
                keys[-1] = value[k]
                del value[k]
                continue
            if isinstance(value[k], list) and not value[k]:
                # When empty list, skip to record it
                del value[k]
                continue
            i = iter_nested_dict(value[k], keys + [k])
            try:
                f = next(i)
                # contain dict
                iters[k] = itertools.chain([f], i)
                del value[k]
            except StopIteration:
                # not contain dict
                units.add(get_unit(k))
        is_value_iterated = False
        for k in key_copy:
            if k in iters:
                yield from iters[k]
                continue
            if not is_value_iterated:
                if 1 == len(units):
                    yield (keys, value)
                    is_value_iterated = True
                else:
                    if k in value:
                        yield (keys + [k], {k: value[k]})
    elif isinstance(value, list):
        a = list(itertools.chain.from_iterable(
            iter_nested_dict(v, keys + [i]) for i, v in enumerate(value)
        ))
        # check combinable; All elements are dicts that don't contain any dict
        # and that have only one value with the same key
        depth = len(keys) + 1
        dict_keys = set()
        for k, v in a:
            if not (len(k) == depth):
                raise ValueError('Data structure not supported.')
            dict_keys |= set(v.keys())
        if 1 == len(dict_keys):
            # combine dicts by replacing the same key and k[-1]
            new_k = next(iter(dict_keys))
            d = {}
            for (k, v) in a:
                d.update({k[-1]: v[new_k]})
            a = [(keys + [new_k], d)]
        yield from a


abbreviations = ['xdr', 'cve', 'os', 'edr']


def make_title(s):
    strings = s if is_container(s) else [s]
    words = list(itertools.chain.from_iterable(
        split_camel_case(word) for s in strings for word in s.split()
    ))
    words = [w.capitalize() if w and w[0].islower() else w for w in words]
    words = [w.upper() if w.lower() in abbreviations else w for w in words]
    return ' '.join(unique_list(words))


def change_sheet_names(dataframes, sheet_names):
    for sheet_name in list(dataframes.keys()):
        if sheet_name in sheet_names:
            new_sheet_name = sheet_names[sheet_name]
            if new_sheet_name in dataframes:
                raise RuntimeError(f'Worksheet {new_sheet_name} already '
                                   'exists.')
            dataframes[new_sheet_name] = dataframes.pop(sheet_name)
            sheet_name = new_sheet_name
        if 32 <= len(sheet_name):
            raise RuntimeError(f"Worksheet name '{sheet_name}' must exceeds "
                               'the 31 character maximum.')


def restore_sheet_names(dataframes, sheet_names):
    for sheet_name in list(dataframes.keys()):
        found = [k for k, v in sheet_names.items() if v == sheet_name]
        if 1 == len(found):
            new_sheet_name = found[0]
            dataframes[new_sheet_name] = dataframes.pop(sheet_name)


def load_config(path):
    with open(path) as f:
        return yaml.safe_load(f)


def load_dataframes(data_source):
    return pandas.read_excel(data_source, engine='openpyxl', sheet_name=None)


datetime_field_name = 'createdDateTime'
datetime_column_name = make_title(datetime_field_name)
general_fields = ['schemaVersion', 'companyId', 'companyName']
general_sheet_name = 'General'
begin_column_name = 'Begin'
end_column_name = 'End'


def to_dataframes(d):
    datetime_value = datetime.datetime.fromisoformat(
        d.pop(datetime_field_name).replace('Z', '')
    )
    dataframes = {}
    dataframes[general_sheet_name] = pandas.DataFrame({
        make_title(k): [d.pop(k)] for k in general_fields if k in d
    })
    dataframes[general_sheet_name].at[0, begin_column_name] = datetime_value
    dataframes[general_sheet_name].at[0, end_column_name] = datetime_value

    for keys, value in iter_nested_dict(d):
        sheet_name = make_title(keys)
        v = pandas.DataFrame({make_title(k): v if is_container(v) else [v]
                              for k, v in value.items()})
        v.insert(0, datetime_column_name, datetime_value)
        dataframes[sheet_name] = v
    return dataframes


def concat_dataframes(*args):
    if args:
        dataframes = args[0]
        general = dataframes[general_sheet_name]
        for arg in args[1:]:
            for sheet_name, dataframe in arg.items():
                if sheet_name in dataframes:
                    if general_sheet_name == sheet_name:
                        begin = dataframe.at[0, begin_column_name]
                        if begin < general.at[0, begin_column_name]:
                            general.at[0, begin_column_name] = begin
                        end = dataframe.at[0, end_column_name]
                        if general.at[0, end_column_name] < end:
                            general.at[0, end_column_name] = end
                        continue
                    dataframes[sheet_name] = pandas.concat([
                        dataframe, dataframes[sheet_name]
                    ])
                else:
                    dataframes[sheet_name] = dataframe
        return dataframes


def to_excel(dataframes, data_source):
    kwargs = {'engine': 'openpyxl', 'mode': 'w'}
    if os.path.isfile(data_source):
        kwargs['mode'] = 'a'
        kwargs['if_sheet_exists'] = 'replace'
    with pandas.ExcelWriter(data_source, **kwargs) as w:
        for sheet_name, dataframe in dataframes.items():
            dataframe.to_excel(w, sheet_name=sheet_name, index=False)


def export(v1, data_source, config):
    sheet_names = config.get('sheet names', {})
    previous_dataframes = []
    if os.path.isfile(data_source):
        dataframes = load_dataframes(data_source)
        restore_sheet_names(dataframes, sheet_names)
        previous_dataframes.append(dataframes)
    dataframes = to_dataframes(v1.get_security_posture())
    dataframes = concat_dataframes(*previous_dataframes, dataframes)
    change_sheet_names(dataframes, sheet_names)
    to_excel(dataframes, data_source)


level_sheet_name = 'Risk Category Level'
level_values = ['low', 'medium', 'high']


def correct_dataframes(dataframes):
    """
    This function correct VisionOne data for presentation

    1. The 'Risk Category Level' has only string values. To make a line chart,
       the value replaced from 'low', 'medium' and 'high' to 0, 1 and 2,
       respectively.

    2. Some sheets have NaN values. To make line charts, the value is replaecd
       to empty string.
    """
    dataframe = dataframes[level_sheet_name]
    for i, row in dataframe.iterrows():
        for column in dataframe.columns:
            if datetime_column_name == column:
                continue
            dataframe.at[i, column] = level_values.index(row[column])

    for dataframe in dataframes.values():
        for i, row in dataframe.iterrows():
            for column in dataframe.columns:
                if datetime_column_name == column:
                    continue
                v = row[column]
                if isinstance(v, numbers.Number) and math.isnan(v):
                    dataframe.at[i, column] = ''


def make_chart_slide(p, title):
    title_only_slide_layout = p.slide_layouts[5]
    slide = p.slides.add_slide(title_only_slide_layout)
    slide.shapes.title.text = title
    lines = len(slide.shapes.title.text) / 30
    if 1 < lines:
        font_size = pptx.util.Pt(44 - 4*lines)
        slide.shapes.title.text_frame.paragraphs[0].font.size = font_size
    return slide


def add_shape_chart_line(p, slide, dataframe):
    chart_data = pptx.chart.data.ChartData()
    for column_name, items in dataframe.items():
        if datetime_column_name == column_name:
            # make category as date axis
            chart_data.categories = [item.to_pydatetime() for item in items]
        else:
            chart_data.add_series(column_name, tuple(item for item in items))
    x = slide.shapes.title.left
    y = slide.shapes.title.top + slide.shapes.title.height
    cx, cy = slide.shapes.title.width, (p.slide_height - y)
    chart = slide.shapes.add_chart(
        pptx.enum.chart.XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.include_in_layout = True
    for s in chart.series:
        s.smooth = False
    chart.font.size = pptx.util.Pt(12)
    return chart


presentation_title = 'Security Posture'


def to_powerpoint(dataframes, presentation):
    p = pptx.Presentation()
    general = dataframes.pop(general_sheet_name)
    presentation_subtitle = os.linesep.join([
        general[make_title(general_fields[2])][0],
        f'{begin_column_name}: {general.at[0, begin_column_name]}',
        f'{end_column_name}: {general.at[0, end_column_name]}'
    ])
    title_slide_layout = p.slide_layouts[0]
    slide = p.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = presentation_title
    slide.placeholders[1].text = presentation_subtitle
    for sheet_name, dataframe in dataframes.items():
        slide = make_chart_slide(p, sheet_name)
        add_shape_chart_line(p, slide, dataframe)
    p.save(presentation)


def report(data_source, presentation, config):
    sheet_names = config.get('sheet names', {})
    dataframes = load_dataframes(data_source)
    restore_sheet_names(dataframes, sheet_names)
    correct_dataframes(dataframes)
    to_powerpoint(dataframes, presentation)


def main(data_source, config, v1_token=None, v1_url=None, presentation=None):
    config = load_config(config)
    if presentation is not None:
        report(data_source, presentation, config)
    else:
        v1 = TmV1Client(v1_token, v1_url)
        export(v1, data_source, config)


if __name__ == '__main__':
    common_parser = argparse.ArgumentParser(add_help=False)
    common_parser.add_argument(
        '-s', '--data-source', default=V1_XLSX_FILENAME,
        help='File name of the XLSX file with the retrieved security metrics. '
             f'Default: `{V1_XLSX_FILENAME}`.'
    )
    common_parser.add_argument(
        '-c', '--config', default=V1_YAML_FILENAME,
        help='File name of the configuration file. Default: '
             f'`{V1_YAML_FILENAME}`.'
    )
    parser = argparse.ArgumentParser(
        description=('Generate a custom report from the Security Posture '
                     'information.'),
        epilog="Refer to examples of the operations 'export' and 'report'.")
    subparsers = parser.add_subparsers(help='')
    export_parser = subparsers.add_parser(
        'export', parents=[common_parser],
        help=('Generate an XLSX file to store the retrieved security metrics '
              'from the Security Posture API.'),
        epilog=(f'Example: python {os.path.basename(__file__)} export '
                f'-s /path/to/file')
    )
    export_parser.add_argument(
        '-t', '--v1-token', default=V1_TOKEN,
        help='Authentication token of your Trend Vision One user account.'
    )
    export_parser.add_argument(
        '-u', '--v1-url', default=TmV1Client.base_url_default,
        help=('URL of the Trend Vision One server for your region.'
              f' The default value is "{TmV1Client.base_url_default}"')
    )
    report_parser = subparsers.add_parser(
        'report', parents=[common_parser],
        help=('Generate a PPTX file to present the security metrics stored in '
              'the XLSX file.'),
        epilog=(f'Example: python {os.path.basename(__file__)} report'
                f'-p /path/to/file')
    )
    report_parser.add_argument(
        '-p', '--presentation', default=V1_PPTX_FILENAME,
        help=('File name of the PPTX file with the metrics. Default: '
              f'`{V1_PPTX_FILENAME}`.'))
    main(**vars(parser.parse_args()))
